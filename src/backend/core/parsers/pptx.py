from pathlib import Path

from pptx import Presentation

from core.models import ParsedDocument
from core.monitoring.logger import get_logger

from .base import BaseParser

logger = get_logger("pptx_parser")


class PPTXParser(BaseParser):
    """Parser for PowerPoint presentations"""

    async def parse(self, file_path: Path, output_dir: Path) -> ParsedDocument:
        """
        Parse a PPTX document and extract text to markdown.

        Creates:
        - content.md: Full presentation as searchable markdown
        - pages/slide_XXX.md: Individual slide files
        """
        logger.info(f"Parsing PPTX: {file_path}")

        prs = Presentation(str(file_path))
        total_word_count = 0
        full_content = []

        # Create pages directory (we call them pages for consistency)
        pages_dir = output_dir / "pages"
        pages_dir.mkdir(parents=True, exist_ok=True)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = [f"# Slide {slide_num}"]

            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_content.append(f"## {slide.shapes.title.text.strip()}")

            # Extract body text from all shapes
            body_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # Avoid duplicating the title
                            if slide.shapes.title and text == slide.shapes.title.text.strip():
                                continue
                            body_texts.append(text)

            if body_texts:
                slide_content.append("\n".join(f"- {t}" for t in body_texts))

            # Extract speaker notes (often contain important context)
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content.append(f"\n**Speaker Notes:**\n{notes_text}")

            slide_text = "\n\n".join(slide_content)

            # Save individual slide
            slide_file = pages_dir / f"slide_{slide_num:03d}.md"
            slide_file.write_text(slide_text, encoding="utf-8")

            # Add to full content
            full_content.append(slide_text)
            total_word_count += self.count_words(slide_text)

        # Save full document
        content_file = output_dir / "content.md"
        full_text = "\n\n---\n\n".join(full_content)
        content_file.write_text(full_text, encoding="utf-8")

        logger.info(f"Parsed PPTX: {len(full_content)} slides, {total_word_count} words")

        return ParsedDocument(
            pageCount=len(full_content),
            wordCount=total_word_count,
            slideCount=len(full_content),
        )
